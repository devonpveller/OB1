"""
openbrain-extract — content extraction sidecar (P5.1).

A stable `POST /extract` contract over a FORMAT → EXTRACTOR REGISTRY, so the
workbench/import pipeline never learns about new formats: a future format
(epub, html, eml, spreadsheet…) is a new registry entry behind the unchanged
contract, not a caller change.

    POST /extract  (multipart: file=<upload>, [content_type=<hint>])
      -> { markdown, title, metadata, pages, images: [{name, b64, mime}] }

⚠ Security (plan §11): this parses UNTRUSTED files/images — a classic RCE
vector. Run unprivileged, no extra network beyond the host STT endpoint, no
shell-outs to convertors that exec embedded scripts. The container is the
sandbox boundary (compose: non-root user, cap_drop, read_only root FS,
restricted networks).

Per-format extractor choice is correctness-first (§12.3): PyMuPDF (PDF),
python-docx / python-pptx (Office), Pillow + Tesseract OCR (images), the local
STT service (audio/video). Each format SHOULD meet its extraction-quality
acceptance gate (text fidelity, tables, headings, image refs) — tracked in the
per-format tests, not enforced here.
"""
from __future__ import annotations

import base64
import io
import os
import subprocess
import tempfile
from typing import Callable

from fastapi import FastAPI, File, Form, HTTPException, UploadFile

app = FastAPI(title="openbrain-extract", version="0.1.0")

# Local STT for audio/video — the ONE allowed egress. Defaults match the
# ai-stack realtime-audio server's one-shot HTTP endpoint:
#   POST http://host.docker.internal:8000/stt  (multipart field "audio_file")
#   -> { "transcript": "...", "language": "...", ... }
# (NOT the OpenAI `/v1/audio/transcriptions` convention.) All overridable.
STT_BASE = os.environ.get("STT_API_BASE", "http://host.docker.internal:8000").rstrip("/")
STT_PATH = os.environ.get("STT_TRANSCRIBE_PATH", "/stt")
STT_FIELD = os.environ.get("STT_FILE_FIELD", "audio_file")
# Fail FAST rather than hanging if the STT service is wrong/slow.
STT_TIMEOUT = float(os.environ.get("STT_TIMEOUT_SEC", "120"))
MAX_BYTES = int(os.environ.get("EXTRACT_MAX_BYTES", str(100 * 1024 * 1024)))


# ── Extraction result shape ────────────────────────────────────────────────
def result(markdown: str, title: str = "", *, pages=None, images=None, **meta):
    return {
        "markdown": markdown or "",
        "title": title or "",
        "metadata": meta or {},
        "pages": pages if pages is not None else [],
        "images": images or [],
    }


# ── Per-format extractors ──────────────────────────────────────────────────
def extract_text(data: bytes, filename: str) -> dict:
    text = data.decode("utf-8", errors="replace")
    title = os.path.splitext(os.path.basename(filename))[0]
    return result(text, title, pages=[text], chars=len(text))


def extract_pdf(data: bytes, filename: str) -> dict:
    import fitz  # PyMuPDF

    doc = fitz.open(stream=data, filetype="pdf")
    pages: list[str] = []
    images: list[dict] = []
    for pno in range(doc.page_count):
        page = doc.load_page(pno)
        pages.append(page.get_text("text"))
        for i, img in enumerate(page.get_images(full=True)):
            try:
                xref = img[0]
                pix = fitz.Pixmap(doc, xref)
                if pix.n - pix.alpha >= 4:  # CMYK / odd → normalize to RGB
                    pix = fitz.Pixmap(fitz.csRGB, pix)
                images.append({
                    "name": f"p{pno + 1}-img{i + 1}.png",
                    "b64": base64.b64encode(pix.tobytes("png")).decode(),
                    "mime": "image/png",
                })
            except Exception:
                continue
    title = (doc.metadata or {}).get("title") or os.path.splitext(os.path.basename(filename))[0]
    md = "\n\n".join(f"<!-- page {i + 1} -->\n{t}" for i, t in enumerate(pages))
    return result(md, title, pages=pages, images=images, page_count=doc.page_count)


def extract_docx(data: bytes, filename: str) -> dict:
    import docx  # python-docx

    d = docx.Document(io.BytesIO(data))
    lines: list[str] = []
    for p in d.paragraphs:
        style = (p.style.name or "").lower() if p.style else ""
        if style.startswith("heading"):
            level = "".join(ch for ch in style if ch.isdigit()) or "1"
            lines.append(f"{'#' * min(int(level), 6)} {p.text}")
        elif p.text.strip():
            lines.append(p.text)
    for table in d.tables:  # simple markdown tables
        for row in table.rows:
            lines.append("| " + " | ".join(c.text.replace("\n", " ") for c in row.cells) + " |")
        lines.append("")
    title = os.path.splitext(os.path.basename(filename))[0]
    return result("\n\n".join(lines), title, pages=["\n\n".join(lines)])


# Legacy OLE binary formats (.doc / .ppt) — python-docx/-pptx only read OOXML.
# catdoc/catppt parse the binary DIRECTLY without executing macros, so they keep
# the sandbox posture (no LibreOffice/headless converter). List-arg subprocess
# (no shell) over a tmpfile in the tmpfs /tmp.
def _convert(cmd: list[str], data: bytes, suffix: str) -> str:
    with tempfile.NamedTemporaryFile(suffix=suffix, delete=False, dir="/tmp") as f:
        f.write(data)
        path = f.name
    try:
        out = subprocess.run(cmd + [path], capture_output=True, timeout=120)
        if out.returncode != 0:
            msg = out.stderr.decode("utf-8", "replace")[:200] or "converter failed"
            raise RuntimeError(msg)
        return out.stdout.decode("utf-8", "replace")
    finally:
        try:
            os.unlink(path)
        except OSError:
            pass


def extract_doc(data: bytes, filename: str) -> dict:
    text = _convert(["catdoc", "-w"], data, ".doc")  # -w: no hard line wrapping
    title = os.path.splitext(os.path.basename(filename))[0]
    return result(text, title, pages=[text])


def extract_ppt(data: bytes, filename: str) -> dict:
    text = _convert(["catppt"], data, ".ppt")
    title = os.path.splitext(os.path.basename(filename))[0]
    return result(text, title, pages=[text])


def extract_pptx(data: bytes, filename: str) -> dict:
    from pptx import Presentation  # python-pptx

    prs = Presentation(io.BytesIO(data))
    slides: list[str] = []
    for n, slide in enumerate(prs.slides, 1):
        chunks = [f"## Slide {n}"]
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                chunks.append(shape.text_frame.text)
        slides.append("\n\n".join(chunks))
    title = os.path.splitext(os.path.basename(filename))[0]
    return result("\n\n".join(slides), title, pages=slides, slide_count=len(slides))


def extract_image(data: bytes, filename: str) -> dict:
    from PIL import Image  # Pillow
    import pytesseract

    img = Image.open(io.BytesIO(data))
    try:
        ocr = pytesseract.image_to_string(img) or ""
    except Exception as e:  # OCR is best-effort; still return the image asset
        ocr = ""
        meta_err = str(e)
    else:
        meta_err = None
    title = os.path.splitext(os.path.basename(filename))[0]
    # A standalone image → a source whose CONTENT is the OCR/caption text (for
    # embedding), with the image itself as an asset.
    images = [{
        "name": os.path.basename(filename),
        "b64": base64.b64encode(data).decode(),
        "mime": Image.MIME.get(img.format or "", "image/png"),
    }]
    return result(ocr, title, images=images, ocr_error=meta_err, width=img.width, height=img.height)


def extract_audio(data: bytes, filename: str) -> dict:
    import httpx

    files = {STT_FIELD: (os.path.basename(filename), data)}
    url = f"{STT_BASE}{STT_PATH}"
    try:
        r = httpx.post(url, files=files, timeout=STT_TIMEOUT)
    except httpx.TimeoutException:
        raise RuntimeError(f"STT timed out after {STT_TIMEOUT}s calling {url} — is the STT service up at STT_API_BASE?")
    except httpx.ConnectError:
        raise RuntimeError(f"STT unreachable at {url} — check STT_API_BASE / host.docker.internal")
    if r.status_code != 200:
        raise RuntimeError(f"STT {r.status_code} at {url}: {r.text[:200]}")
    body = r.json()
    # The realtime-audio server returns `transcript`; tolerate `text` too.
    text = body.get("transcript") or body.get("text") or ""
    title = os.path.splitext(os.path.basename(filename))[0]
    return result(text, title, pages=[text], transcribed=True, stt_language=body.get("language"))


# ── Registry (extension → handler). A new format = one entry here. ──────────
REGISTRY: dict[str, Callable[[bytes, str], dict]] = {
    "txt": extract_text, "md": extract_text, "markdown": extract_text,
    "pdf": extract_pdf,
    "docx": extract_docx, "doc": extract_doc,
    "pptx": extract_pptx, "ppt": extract_ppt,
    "png": extract_image, "jpg": extract_image, "jpeg": extract_image,
    "gif": extract_image, "webp": extract_image, "bmp": extract_image, "tiff": extract_image,
    "mp3": extract_audio, "wav": extract_audio, "m4a": extract_audio,
    "mp4": extract_audio, "mov": extract_audio, "webm": extract_audio, "ogg": extract_audio,
}

# content_type the resulting `sources` row should carry per extension group.
CONTENT_TYPE = {
    "pdf": "pdf", "docx": "docx", "pptx": "pptx",
    "doc": "docx", "ppt": "pptx",  # legacy → same document kind (source_format keeps .doc/.ppt)
    "txt": "txt", "md": "md", "markdown": "md",
    "png": "image", "jpg": "image", "jpeg": "image", "gif": "image",
    "webp": "image", "bmp": "image", "tiff": "image",
    "mp3": "audio", "wav": "audio", "m4a": "audio",
    "mp4": "audio", "mov": "audio", "webm": "audio", "ogg": "audio",
}


@app.get("/health")
def health():
    return {"ok": True, "service": "openbrain-extract", "formats": sorted(REGISTRY.keys())}


@app.post("/extract")
async def extract(file: UploadFile = File(...), content_type: str = Form(default="")):
    data = await file.read()
    if not data:
        raise HTTPException(400, "empty upload")
    if len(data) > MAX_BYTES:
        raise HTTPException(413, f"file exceeds {MAX_BYTES} bytes")
    ext = (os.path.splitext(file.filename or "")[1].lstrip(".") or content_type).lower()
    handler = REGISTRY.get(ext)
    if not handler:
        raise HTTPException(415, f"unsupported format: {ext!r} (supported: {sorted(REGISTRY)})")
    try:
        out = handler(data, file.filename or f"upload.{ext}")
    except HTTPException:
        raise
    except Exception as e:  # corrupt/unparseable file → clear error (gate)
        raise HTTPException(422, f"extraction failed for {ext}: {e}")
    out["metadata"]["content_type"] = CONTENT_TYPE.get(ext, "web_article")
    out["metadata"]["source_format"] = ext
    return out
